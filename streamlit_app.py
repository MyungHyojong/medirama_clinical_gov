import streamlit as st
import pandas as pd
import requests
import openai
import os
import io
import zipfile
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from openai import OpenAI

# ========== 환경설정 ========== #
# [주의] 실제 사용 시에는 아래처럼 직접 키를 하드코딩하지 말고,
#       환경변수(os.getenv)나 st.secrets 등을 사용하세요.
openai.api_key = st.secrets["openai"]["api_key"]
client = OpenAI(api_key=openai.api_key)

# --------------------- 세션 스테이트 초기화 --------------------- #
if "processed_df" not in st.session_state:
    st.session_state["processed_df"] = None
if "html_contents" not in st.session_state:
    st.session_state["html_contents"] = {}
if "gemini_responses" not in st.session_state:
    st.session_state["gemini_responses"] = {}

# ---------------------------------------------------------------------------- #
#                             헬퍼 함수 (Utilities)                            #
# ---------------------------------------------------------------------------- #

def find_common_substring(s1, s2):
    if not s1 or not s2:
        return ""
    max_len = min(len(s1), len(s2))
    longest_common_substr = ""
    for length in range(max_len, 0, -1):
        for start1 in range(len(s1) - length + 1):
            substr1 = s1[start1:start1 + length]
            for start2 in range(len(s2) - length + 1):
                substr2 = s2[start2:start2 + length]
                if substr1 == substr2:
                    return substr1
    return ""

def extract_last_in_brackets(text):
    import re
    matches = re.findall(r"\(([^)]+)\)", text)
    if matches:
        return matches[-1]
    return None

def highlight_substring(text, substring):
    if not substring:
        return text
    return text.replace(substring, f"<mark>{substring}</mark>")

def truncate_text(x, max_len=20):
    if isinstance(x, str) and len(x) > max_len:
        return x[:max_len] + "..."
    return x

def read_excel_bytes(df):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name="Sheet1")
    return output.getvalue()

def zip_html_files(html_dict):
    output = io.BytesIO()
    with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zf:
        for fname, html_str in html_dict.items():
            if html_str is None:
                continue
            zf.writestr(fname, html_str)
    return output.getvalue()

def create_ppt_from_dfs(
    df_list_dict,       # {"category": DataFrame, ...}
    common_title,       # 공통 슬라이드 제목
    selected_columns,   # PPT에 표시할 컬럼 목록
    rows_per_slide=6,
    col_widths=None     # { "컬럼명": 퍼센트(0~100), ... }
):
    prs = Presentation()
    slide_width_in_EMU = prs.slide_width  # 기본 10인치 = 9,144,000 EMU

    title_slide_layout = prs.slide_layouts[5]  # "Title and Content"

    margin_in_inches = 0.5
    margin_in_EMU = int(margin_in_inches * 914400)
    table_width_in_EMU = slide_width_in_EMU - (2 * margin_in_EMU)

    for category_name, df_data in df_list_dict.items():
        if df_data.empty:
            continue

        df_data = df_data[selected_columns]
        total_rows = len(df_data)
        num_splits = (total_rows + rows_per_slide - 1) // rows_per_slide

        for idx in range(num_splits):
            slide = prs.slides.add_slide(title_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            if not title_shape:
                title_shape = shapes.add_textbox(
                    margin_in_EMU, margin_in_EMU,
                    slide_width_in_EMU - 2*margin_in_EMU,
                    int(0.5 * 914400)
                )

            slide_title = f"{common_title} - {category_name}"
            if num_splits > 1:
                slide_title += f" (page {idx+1})"
            title_shape.text = slide_title

            left = margin_in_EMU
            top = int(1.5 * 914400)
            table_height_in_EMU = int(5 * 914400)

            df_chunk = df_data.iloc[idx*rows_per_slide : (idx+1)*rows_per_slide]
            rows = len(df_chunk) + 1
            cols = len(selected_columns)

            table = shapes.add_table(rows, cols, left, top, table_width_in_EMU, table_height_in_EMU).table

            # 헤더
            for c, col_name in enumerate(selected_columns):
                cell = table.cell(0, c)
                cell.text = col_name
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.runs[0]
                run.font.bold = True
                run.font.size = Pt(12)

            # 실제 데이터
            for r in range(len(df_chunk)):
                for c, col_name in enumerate(selected_columns):
                    val = df_chunk.iloc[r][col_name]
                    val_str = "" if pd.isnull(val) else str(val)
                    table.cell(r+1, c).text = val_str

            # 열 너비 (%)
            for c, col_name in enumerate(selected_columns):
                if col_widths and (col_name in col_widths):
                    col_percent = col_widths[col_name] / 100.0
                else:
                    col_percent = 1.0 / cols
                col_emu_width = int(table_width_in_EMU * col_percent)
                table.columns[c].width = col_emu_width

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def main():
    st.title("Clinical Trial Data Processor")

    tab1, tab2 = st.tabs(["Upload and Process", "Select Rows and Generate PPT"])

    with tab1:
        st.header("1) CSV 업로드 후 처리하기")

        uploaded_file = st.file_uploader("CSV 파일 업로드", type=["csv"])
        process_button = st.button("Process Data")

        if process_button and uploaded_file:
            df = pd.read_csv(uploaded_file, encoding="utf-8")
            if "NCT Number" not in df.columns:
                st.error("업로드된 CSV 파일에 'NCT Number' 컬럼이 없습니다.")
                return

            nct_list = df["NCT Number"].tolist()

            tests = []
            reasons = []
            explanations = []
            genes_list = []
            confidence_scores = []
            experimental_1s = []
            experimental_2s = []
            control_1s = []
            control_2s = []
            study_names = []
            start_dates = []
            primary_completion_dates = []
            completion_dates = []

            # 추가: Official/Brief Title을 저장할 리스트
            official_titles = []
            brief_titles = []

            html_contents = {}
            gemini_responses = {}

            progress_bar = st.progress(0)
            total_ncts = len(nct_list)

            for i, nct_id in enumerate(nct_list):
                try:
                    url = f"https://clinicaltrials.gov/api/v2/studies/{nct_id}"
                    params = {
                        'format': 'json',
                        'markupFormat': 'markdown',
                        'fields': 'NCTId,BriefTitle,ConditionsModule,EligibilityModule,OfficialTitle,ArmsInterventionsModule,StartDate,PrimaryCompletionDate,CompletionDate'
                    }
                    response = requests.get(url, params=params)

                    if response.status_code != 200:
                        tests.append(None)
                        reasons.append(None)
                        explanations.append(None)
                        genes_list.append(None)
                        confidence_scores.append(None)
                        experimental_1s.append(None)
                        experimental_2s.append(None)
                        control_1s.append(None)
                        control_2s.append(None)
                        study_names.append(None)
                        start_dates.append(None)
                        primary_completion_dates.append(None)
                        completion_dates.append(None)

                        # Official/Brief Title도 None
                        official_titles.append(None)
                        brief_titles.append(None)

                        html_contents[f"{nct_id}.html"] = None
                        gemini_responses[nct_id] = None
                        continue

                    study_data = response.json()
                    arms_data = study_data.get('protocolSection', {}).get('armsInterventionsModule', {}).get('armGroups')

                    exp1, exp2, ctr1, ctr2 = None, None, None, None
                    if arms_data:
                        experimental_labels_drugs = []
                        experimental_descriptions = []
                        control_labels_drugs = []
                        control_descriptions = []

                        for arm in arms_data:
                            label = arm.get('label', '')
                            arm_type = arm.get('type', '')
                            description = arm.get('description', '')

                            other_keys = []
                            for k, v in arm.items():
                                if k not in ['label', 'type', 'description']:
                                    if isinstance(v, str):
                                        other_keys.append(v)
                                    elif isinstance(v, list):
                                        other_keys.append('/'.join(map(str, v)))
                            drug_info = "\n".join(other_keys) if other_keys else "No drug info"

                            if 'EXPERIMENTAL' in arm_type.upper():
                                experimental_labels_drugs.append(f"{label}\n{drug_info}")
                                experimental_descriptions.append(description)
                            else:
                                control_labels_drugs.append(f"{label}\n{drug_info}")
                                control_descriptions.append(description)

                        exp1 = "\n\n".join(experimental_labels_drugs)
                        exp2 = "\n\n".join(experimental_descriptions)
                        ctr1 = "\n\n".join(control_labels_drugs)
                        ctr2 = "\n\n".join(control_descriptions)

                    experimental_1s.append(exp1)
                    experimental_2s.append(exp2)
                    control_1s.append(ctr1)
                    control_2s.append(ctr2)

                    # GPT
                    experiment_description = str(study_data)
                    question = (
                        "Please provide an answer in the following format based on the provided experiment description,\n"
                        "1. test: Determine which of the following categories the cancer treatment experiment belongs to: first line, second line, third line, neoadjuvant, adjuvant or unclear. Only print out the test type without writinig any sentences.\n"
                        "2. reason: write the exact specific part of the eligibilityCriteria in Experiment description below that supports your answer in 1 without changing a single part\n"
                        "3. explanations: explain specifically why you chose your answer in 1\n"
                        "4. genes: mutations, expressions associated in study such as KRAS, EGFR, MET, ALK, CEACAM5, STK11, KEAP1\n"
                        "5. Confidence Score: Based on the probability of incorrectly guessing the type of test, if you are uncertain about your decision, write 'uncertain'. else write 'certain'.\n\n"
                        f"Experiment description below:\n{experiment_description}"
                    )

                    try:
                        response = client.chat.completions.create(
                            model="gpt-3.5-turbo",
                            messages=[
                                {
                                    "role": "system",
                                    "content": "You are a helpful assistant that analyzes clinical trial descriptions."
                                },
                                {
                                    "role": "user",
                                    "content": question
                                }
                            ]
                        )
                        response_text = response.choices[0].message.content
                    except Exception:
                        response_text = None

                    test_val, reason_val, explanation_val = None, None, None
                    genes_val, confidence_val = None, None

                    if response_text:
                        gemini_responses[nct_id] = response_text
                        lines = response_text.split("\n")
                        for line in lines:
                            line_stripped = line.strip().lower()
                            if line_stripped.startswith("1."):
                                chunk = line.strip().split("1.", 1)[-1]
                                test_val = chunk.replace("test:", "").strip().lower()
                            elif line_stripped.startswith("2."):
                                reason_val = line.strip().split("2.", 1)[-1].strip()
                            elif line_stripped.startswith("3."):
                                explanation_val = line.strip().split("3.", 1)[-1].strip()
                            elif line_stripped.startswith("4."):
                                genes_val = line.strip().split("4.", 1)[-1].strip()
                            elif line_stripped.startswith("5."):
                                confidence_val = line.strip().split("5.", 1)[-1].strip()
                    else:
                        gemini_responses[nct_id] = None

                    tests.append(test_val)
                    reasons.append(reason_val)
                    explanations.append(explanation_val)
                    genes_list.append(genes_val)
                    confidence_scores.append(confidence_val)

                    # 공식 타이틀, brief title, conditions, eligibility
                    identification_module = study_data.get('protocolSection', {}).get('identificationModule', {})
                    o_title = identification_module.get('officialTitle', "NA")
                    b_title = identification_module.get('briefTitle', "NA")

                    # 리스트에 저장
                    official_titles.append(o_title)
                    brief_titles.append(b_title)

                    conditions_module = study_data.get('protocolSection', {}).get('conditionsModule', {})
                    conditions_val = conditions_module.get('conditions', [])
                    conditions = ", ".join(conditions_val) if conditions_val else "NA"

                    eligibility_module = study_data.get('protocolSection', {}).get('eligibilityModule', {})
                    eligibility_criteria = eligibility_module.get('eligibilityCriteria', "")

                    found_name_in_brackets = extract_last_in_brackets(o_title) or "NA"
                    study_names.append(found_name_in_brackets)

                    status_module = study_data.get('protocolSection', {}).get('statusModule', {})
                    start_date_info = status_module.get('startDateStruct')
                    primary_completion_info = status_module.get('primaryCompletionDateStruct')
                    completion_info = status_module.get('completionDateStruct')

                    def parse_date_info(d):
                        if not d:
                            return "NA"
                        return str(d)

                    start_dates.append(parse_date_info(start_date_info))
                    primary_completion_dates.append(parse_date_info(primary_completion_info))
                    completion_dates.append(parse_date_info(completion_info))

                    reason_parts = find_common_substring(reason_val if reason_val else "", eligibility_criteria)
                    highlighted_criteria = highlight_substring(eligibility_criteria, reason_parts)
                    highlighted_criteria = highlighted_criteria.replace("\n", "<br>")

                    arms_html = (
                        "<table border='1' style='border-collapse:collapse'>"
                        "<tr><th>label</th><th>type</th><th>description</th><th>others</th></tr>"
                    )
                    if arms_data:
                        for arm in arms_data:
                            label = arm.get('label', '')
                            arm_type = arm.get('type', '')
                            desc = arm.get('description', '')
                            others = []
                            for k, v in arm.items():
                                if k not in ['label', 'type', 'description']:
                                    others.append(f"{k}:{v}")
                            arms_html += f"<tr><td>{label}</td><td>{arm_type}</td><td>{desc}</td><td>{'<br>'.join(others)}</td></tr>"
                    arms_html += "</table>"

                    markdown_content = (
                        f"<h1>Response</h1>"
                        f"<strong>Test:</strong> {test_val}<br><br>"
                        f"<strong>Explanation:</strong> {explanation_val}<br><br>"
                        f"<strong>Genes:</strong> {genes_val}<br><br>"
                        f"<strong>Confidence Score:</strong> {confidence_val}<br><br>"
                        f"<h1>Study Plan</h1>"
                        f"{arms_html}"
                        f"<h1>Study Data</h1>"
                        f"<strong>Official Title:</strong> {o_title}<br><br>"
                        f"<strong>Brief Title:</strong> {b_title}<br><br>"
                        f"<strong>Conditions:</strong> {conditions}<br><br>"
                        f"<strong>Eligibility Module:</strong><br>{highlighted_criteria}<br><br>"
                        f"<strong>Highlighted Reason from Response:</strong><br>"
                        f"{('<mark>' + reason_parts + '</mark>') if reason_parts else 'N/A'}"
                    )
                    html_contents[f"{nct_id}.html"] = markdown_content

                except Exception:
                    tests.append(None)
                    reasons.append(None)
                    explanations.append(None)
                    genes_list.append(None)
                    confidence_scores.append(None)
                    experimental_1s.append(None)
                    experimental_2s.append(None)
                    control_1s.append(None)
                    control_2s.append(None)
                    study_names.append(None)
                    start_dates.append(None)
                    primary_completion_dates.append(None)
                    completion_dates.append(None)
                    official_titles.append(None)
                    brief_titles.append(None)
                    html_contents[f"{nct_id}.html"] = None
                    gemini_responses[nct_id] = None

                progress_bar.progress((i+1)/total_ncts)

            # df에 새로운 컬럼들 추가
            df["official_title"] = official_titles
            df["brief_title"] = brief_titles
            df["TestLines"] = tests
            df["Reason"] = reasons
            df["Explanation"] = explanations
            df["Genes"] = genes_list
            df["Confidence_Score"] = confidence_scores
            df["experimental_data"] = experimental_1s
            df["experimental_description"] = experimental_2s
            df["control_datas"] = control_1s
            df["control_description"] = control_2s
            df["study_name"] = study_names
            df["start_date"] = start_dates
            df["primary_completion_date"] = primary_completion_dates
            df["completion_date"] = completion_dates

            st.session_state["processed_df"] = df
            st.session_state["html_contents"] = html_contents
            st.session_state["gemini_responses"] = gemini_responses

            st.success("Data processing complete.")

        # 처리된 데이터 표시
        if st.session_state["processed_df"] is not None:
            df_display = st.session_state["processed_df"].copy()
            for c in df_display.columns:
                df_display[c] = df_display[c].apply(lambda x: truncate_text(x, 20))
            st.dataframe(df_display)

            # Excel 다운로드
            excel_bytes = read_excel_bytes(st.session_state["processed_df"])
            st.download_button(
                label="Download Updated Data (Excel)",
                data=excel_bytes,
                file_name="updated_ctg_studies.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )

            # HTML(zip) 다운로드
            zip_bytes = zip_html_files(st.session_state["html_contents"])
            st.download_button(
                label="Download HTML Files",
                data=zip_bytes,
                file_name="html_files.zip",
                mime="application/zip"
            )

            # GPT Raw Response
            st.subheader("GPT Response")
            valid_responses = [v for v in st.session_state["gemini_responses"].values() if v]
            if not valid_responses:
                st.write("No responses available.")
            else:
                st.write("\n\n".join(valid_responses))
        else:
            st.info("CSV 파일을 업로드한 뒤, 'Process Data' 버튼을 눌러주세요.")

    # ========================================================================== 
    # TAB 2: PPT 생성 (CSV/XLSX 업로드, % 단위 열 너비 설정)
    # ========================================================================== 
    with tab2:
        st.header("2) 파일 업로드 후 PPT 생성")

        ppt_file = st.file_uploader("CSV / XLSX 파일 업로드 (PPT용)", type=["csv", "xlsx"], key="ppt_file")
        common_title = st.text_input("공통 슬라이드 제목", value="", placeholder="슬라이드에 공통으로 들어갈 제목 입력")

        if ppt_file is not None:
            import os
            ext = os.path.splitext(ppt_file.name)[-1].lower()

            if ext == ".csv":
                df_ppt = pd.read_csv(ppt_file, encoding="utf-8")
            else:
                df_ppt = pd.read_excel(ppt_file)

            all_cols = df_ppt.columns.tolist()
            grouping_col = st.selectbox("그룹화할 컬럼 선택 (선택 사항)", ["None"] + all_cols)
            rows_per_slide = st.slider("슬라이드 당 행 개수", 1, 20, 6)

            selected_columns = st.multiselect(
                "PPT에 포함할 컬럼 선택",
                all_cols,
                default=all_cols[:3]
            )

            col_widths = {}
            if selected_columns:
                for col in selected_columns:
                    col_widths[col] = st.slider(
                        f"{col} 열 너비 (%)",
                        min_value=0,
                        max_value=100,
                        value=25,  # 기본값 예시
                        step=1
                    )

            generate_ppt = st.button("Generate PPT")
            if generate_ppt and selected_columns:
                if grouping_col != "None":
                    categories = df_ppt[grouping_col].unique()
                    df_list_dict = {}
                    for cat in categories:
                        cat_df = df_ppt[df_ppt[grouping_col] == cat]
                        df_list_dict[str(cat)] = cat_df
                else:
                    df_list_dict = {"Data": df_ppt}

                with st.spinner("Generating PPT..."):
                    ppt_io = create_ppt_from_dfs(
                        df_list_dict=df_list_dict,
                        common_title=common_title,
                        selected_columns=selected_columns,
                        rows_per_slide=rows_per_slide,
                        col_widths=col_widths
                    )
                st.success("PPT generation complete.")

                st.download_button(
                    label="Download PPT",
                    data=ppt_io,
                    file_name="selected_trials.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

                st.subheader("Preview of Dataframes per Category")
                for cat_name, dataf in df_list_dict.items():
                    if not dataf.empty:
                        preview_df = dataf[selected_columns].head(rows_per_slide)
                        st.write(f"**Category: {cat_name}**")
                        st.dataframe(preview_df)
                    else:
                        st.write(f"**Category: {cat_name}** (Empty)")

        else:
            st.info("CSV/XLSX 파일을 업로드하세요.")


if __name__ == "__main__":
    main()
